"""Generate Friday review deck — clean PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from lxml import etree
from pptx.oxml.ns import qn

# ── palette ──────────────────────────────────────────────
BG       = RGBColor(0x08, 0x08, 0x0E)
BG2      = RGBColor(0x0E, 0x0E, 0x18)
BG3      = RGBColor(0x12, 0x12, 0x1E)
WHITE    = RGBColor(0xF0, 0xF0, 0xF6)
LGREY    = RGBColor(0xA0, 0xA0, 0xB8)
GREY     = RGBColor(0x60, 0x60, 0x78)
DGREY    = RGBColor(0x28, 0x28, 0x3A)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
DKBLUE   = RGBColor(0x06, 0x0E, 0x1A)
BRBLUE   = RGBColor(0x38, 0xBD, 0xF8)
BDRBLUE  = RGBColor(0x1A, 0x34, 0x54)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
DKGREEN  = RGBColor(0x06, 0x12, 0x10)
BDRGREEN = RGBColor(0x10, 0x38, 0x20)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
DKVIOLET = RGBColor(0x0C, 0x08, 0x18)
BDRVIOLET= RGBColor(0x2E, 0x1A, 0x4A)
CYAN     = RGBColor(0x2D, 0xD4, 0xBF)
DKCYAN   = RGBColor(0x05, 0x12, 0x10)
BDRCYAN  = RGBColor(0x0F, 0x30, 0x28)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
DKAMBER  = RGBColor(0x14, 0x10, 0x06)
BDRAMBER = RGBColor(0x3A, 0x28, 0x0A)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
DKORANGE = RGBColor(0x14, 0x0A, 0x06)
BDRORANGE= RGBColor(0x3A, 0x18, 0x0A)
RED      = RGBColor(0xEF, 0x44, 0x44)
DKRED    = RGBColor(0x14, 0x06, 0x06)
BDRRED   = RGBColor(0x3A, 0x10, 0x10)
BORDER   = RGBColor(0x20, 0x20, 0x30)
SOFTGRN  = RGBColor(0x86, 0xEF, 0xAC)
DKSFGRN  = RGBColor(0x0C, 0x12, 0x08)
BDRSFGRN = RGBColor(0x1E, 0x32, 0x10)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── primitives ───────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl


def rect(sl, x, y, w, h, fill=BG2, border=BORDER, bw=Pt(0.75)):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = bw
    return s


def txt(sl, text, x, y, w, h, size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def multirun(sl, runs, x, y, w, h, align=PP_ALIGN.LEFT):
    """runs = [(text, size, bold, color), ...]  all on one paragraph"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def arrow_right(sl, x1, y, x2, col=DGREY, lw=Pt(1.25)):
    """Horizontal arrow from (x1,y) to (x2,y)."""
    c = sl.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    c.line.color.rgb = col; c.line.width = lw
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'triangle'); he.set('w', 'sm'); he.set('len', 'sm')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'none')
    return c


def arrow_down(sl, x, y1, y2, col=DGREY, lw=Pt(1.25)):
    """Vertical arrow from (x,y1) to (x,y2)."""
    c = sl.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x), Inches(y1), Inches(x), Inches(y2),
    )
    c.line.color.rgb = col; c.line.width = lw
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'triangle'); he.set('w', 'sm'); he.set('len', 'sm')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'none')
    return c


def flow_box(sl, x, y, w, h, title, steps,
             title_color=WHITE, fill=BG2, border=BORDER):
    """Box with title + divider + step list."""
    rect(sl, x, y, w, h, fill=fill, border=border)
    txt(sl, title, x+0.14, y+0.10, w-0.28, 0.30,
        size=9.5, bold=True, color=title_color)
    # thin divider
    div = sl.shapes.add_shape(1,
        Inches(x+0.12), Inches(y+0.44),
        Inches(w-0.24), Emu(6000))
    div.fill.solid(); div.fill.fore_color.rgb = border
    div.line.fill.background()
    for i, step in enumerate(steps):
        txt(sl, step, x+0.14, y+0.52 + i*0.30,
            w-0.28, 0.27, size=8.5, color=LGREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — The Idea
# ═══════════════════════════════════════════════════════════
sl1 = add_slide()

txt(sl1, "HEBBIA MATRIX  ·  PROOF OF CONCEPT",
    0.55, 0.42, 9, 0.28, size=9, bold=True, color=BLUE)

multirun(sl1, [
    ("Turning documents into a\n", 34, True, WHITE),
    ("structured, queryable grid", 34, True, BLUE),
], 0.55, 0.78, 9.2, 1.3)

txt(sl1,
    "Hebbia Matrix lets analysts define questions as columns and run them "
    "across every document simultaneously — each cell is a cited, verified answer "
    "traced back to the exact source page.",
    0.55, 2.22, 7.8, 0.85, size=12.5, color=GREY)

# three cards
cards = [
    ("ROWS",    "Documents",        "10-Ks, filings, transcripts — uploaded once, cached permanently", BLUE),
    ("COLUMNS", "Questions",        "Any analyst question. Edit a column and stale cells rerun automatically", BLUE),
    ("CELLS",   "Cited Answers",    "Every claim cites a page and passage. Full reasoning trace on click", GREEN),
]
for i, (tag, title, body, col) in enumerate(cards):
    cx = 0.55 + i * 2.82
    rect(sl1, cx, 3.28, 2.6, 1.52, fill=BG2, border=BORDER)
    txt(sl1, tag,   cx+0.18, 3.38, 2.2, 0.24, size=8,  bold=True,  color=col)
    txt(sl1, title, cx+0.18, 3.66, 2.2, 0.32, size=13, bold=True,  color=WHITE)
    txt(sl1, body,  cx+0.18, 4.02, 2.2, 0.65, size=10, color=GREY)

# mini grid — right side
GX, GY = 9.55, 0.9
rect(sl1, GX, GY, 3.55, 3.2, fill=BG3, border=BORDER)

# header
rect(sl1, GX, GY, 3.55, 0.34, fill=BG2, border=BORDER)
for hx, ht in [(0.10, "Document"), (1.22, "Revenue"), (2.38, "Auditor")]:
    txt(sl1, ht, GX+hx, GY+0.07, 1.1, 0.22, size=7.5, bold=True, color=BLUE)

rows = [
    ("Apple 10-K",  "-2.8%  p.42",  "EY  p.91",   GREEN,  GREEN),
    ("NVIDIA 10-K", "Generating…",  "Pending",     BRBLUE, GREY),
    ("Netflix 10-K","Pending",      "Pending",      GREY,   GREY),
]
for ri, (doc, rev, aud, rc, ac) in enumerate(rows):
    ry = GY + 0.34 + ri * 0.94
    rect(sl1, GX, ry, 3.55, 0.94,
         fill=BG if ri % 2 == 0 else BG3, border=BORDER)
    txt(sl1, doc, GX+0.10, ry+0.30, 1.1, 0.28, size=8,   color=LGREY)
    txt(sl1, rev, GX+1.22, ry+0.30, 1.1, 0.28, size=8,   bold=True, color=rc)
    txt(sl1, aud, GX+2.38, ry+0.30, 0.9, 0.28, size=8,   color=ac)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — What Was Built / Where It Fits
# ═══════════════════════════════════════════════════════════
sl2 = add_slide()

txt(sl2, "WHAT WAS BUILT  ·  WHERE IT GOES",
    0.55, 0.42, 9, 0.28, size=9, bold=True, color=VIOLET)

multirun(sl2, [
    ("What was built  ", 28, True, WHITE),
    ("and where it fits", 28, True, VIOLET),
], 0.55, 0.78, 10, 0.85)

# pipeline strip
rect(sl2, 0.55, 1.80, 12.2, 0.62, fill=RGBColor(0x09,0x09,0x14), border=BORDER)
stages = ["Vision Ingest", "Wiki Builder", "Vector Index", "ISD Agent Loop", "Streaming UI"]
stage_colors = [BRBLUE, VIOLET, CYAN, GREEN, BLUE]
sw = 12.2 / len(stages)
for i, (stage, col) in enumerate(zip(stages, stage_colors)):
    sx = 0.55 + i * sw
    txt(sl2, stage, sx + sw*0.5 - 0.6, 1.88, 1.2, 0.22,
        size=8.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl2, ["Batch 5 pages/call","Per-section LLM","Chunks + embeddings","Decompose→verify","React + SSE"][i],
        sx + sw*0.5 - 0.7, 2.12, 1.4, 0.22,
        size=7.5, color=GREY, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl2, "→", sx + sw - 0.12, 1.97, 0.25, 0.22,
            size=10, color=DGREY, align=PP_ALIGN.CENTER)

# three info columns
cols3 = [
    {
        "tag": "WHAT WAS BUILT", "col": GREEN, "fill": DKGREEN, "bdr": BDRGREEN, "x": 0.55,
        "items": [
            "GPT-4.1 Vision ingest — PDF pages as images",
            "Per-document wiki with entities, metrics, claims",
            "Three swappable retriever modes",
            "ISD loop: decompose, gather, draft, verify, revise",
            "Streaming React UI with PDF viewer and citations",
            "FinanceBench harness to benchmark retrieval quality",
        ],
    },
    {
        "tag": "LIMITATIONS", "col": RED, "fill": DKRED, "bdr": BDRRED, "x": 4.57,
        "items": [
            "Local only — no auth, single user",
            "Wiki build cost ~$2–5 per 10-K filing",
            "Page-level citations only, no sub-page highlights",
            "Charts and images inside PDFs not extracted",
            "No document versioning across uploads",
        ],
    },
    {
        "tag": "TO SCALE", "col": AMBER, "fill": DKAMBER, "bdr": BDRAMBER, "x": 8.60,
        "items": [
            "Cloud storage + managed vector database",
            "Multi-tenant auth with workspace isolation",
            "Async job workers with priority queues",
            "Wiki cost shared via content-addressable cache",
            "Sub-page citation via layout extraction",
        ],
    },
]
for c in cols3:
    rect(sl2, c["x"], 2.62, 3.78, 2.52, fill=c["fill"], border=c["bdr"])
    txt(sl2, c["tag"], c["x"]+0.18, 2.72, 3.4, 0.26, size=8, bold=True, color=c["col"])
    for ii, item in enumerate(c["items"]):
        txt(sl2, item, c["x"]+0.22, 3.05 + ii*0.36, 3.38, 0.30, size=10, color=LGREY)

# bigger play band
rect(sl2, 0.55, 5.38, 12.2, 1.62, fill=DKVIOLET, border=BDRVIOLET)
txt(sl2, "THE BIGGER PLAY", 0.78, 5.50, 4, 0.26, size=8, bold=True, color=VIOLET)
plays = [
    ("Matrix output builds a custom structured database from raw documents.", 0.78, 5.82),
    ("That database becomes the backend for an MCP chat interface.", 0.78, 6.20),
    ("Analysts query their own knowledge base conversationally — every answer traced to a source page.", 0.78, 6.58),
]
for ptxt, px, py in plays:
    txt(sl2, ptxt, px, py, 11.7, 0.30, size=11.5, color=RGBColor(0xC4,0xB5,0xFD))


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ═══════════════════════════════════════════════════════════
sl3 = add_slide()

txt(sl3, "ARCHITECTURE", 0.55, 0.36, 5, 0.26, size=9, bold=True, color=GREEN)
multirun(sl3, [
    ("System flow  ", 24, True, WHITE),
    ("— how a PDF becomes a cited answer", 24, False, GREY),
], 0.55, 0.68, 12, 0.70)

# ── OFFLINE BAND ──────────────────────────────────────────
rect(sl3, 0.42, 1.50, 12.50, 1.90, fill=RGBColor(0x0A,0x0A,0x14), border=BORDER)
txt(sl3, "OFFLINE  ·  ONCE PER DOCUMENT  ·  CACHED",
    0.60, 1.55, 6, 0.24, size=7.5, bold=True, color=DGREY)

offline_boxes = [
    ("PDF",              [],                                    WHITE,   BG2,      BORDER),
    ("Vision Ingest",   ["Render pages to PNG","Batch 5 images per call","GPT-4.1 extracts markdown"], BRBLUE, DKBLUE, BDRBLUE),
    ("Wiki Builder",    ["One LLM call per section","Entities, claims, metrics","Doc-level rollup"],    VIOLET, DKVIOLET,BDRVIOLET),
    ("Vector Index",    ["500-token chunks","Embed via Azure / BGE","Store in LanceDB"],                CYAN,   DKCYAN,  BDRCYAN),
]
BW, BH, BGX, BGY = 2.62, 1.46, 0.52, 1.72
BGAP = 0.46
for i, (title, steps, col, fill, bdr) in enumerate(offline_boxes):
    bx = BGX + i * (BW + BGAP)
    flow_box(sl3, bx, BGY, BW, BH, title, steps,
             title_color=col, fill=fill, border=bdr)
    if i < 3:
        arrow_right(sl3, bx + BW, BGY + BH/2,
                    bx + BW + BGAP, col=DGREY)

# ── RETRIEVER BAND ───────────────────────────────────────
rect(sl3, 0.42, 3.54, 12.50, 2.22, fill=RGBColor(0x0A,0x0C,0x10), border=BORDER)
txt(sl3, "RETRIEVER  ·  PLUGGABLE AT RUNTIME — SAME INTERFACE, SWAPPED VIA CONFIG",
    0.60, 3.59, 10, 0.24, size=7.5, bold=True, color=DGREY)

retrievers = [
    (
        "Naive Retriever",
        [
            "1.  Embed query → dense vector",
            "2.  ANN search in LanceDB",
            "3.  Return top-k chunks by cosine score",
        ],
        BRBLUE, DKBLUE, BDRBLUE,
    ),
    (
        "ISD Retriever",
        [
            "1.  LLM: decompose query → sub-queries + target sections",
            "2.  Vector search per sub-query, boost target sections",
            "3.  Collect top-30 candidate chunks, deduplicate",
            "4.  LLM: attention-score each chunk (0.0 – 1.0)",
            "5.  Return top-k by attention score",
        ],
        VIOLET, DKVIOLET, BDRVIOLET,
    ),
    (
        "Wiki Retriever",
        [
            "1.  Search pre-built wiki metrics and claims",
            "2.  Match section via questions_answered index",
            "3.  Drill into that section's chunks",
            "4.  Vector fallback for anything not found above",
        ],
        CYAN, DKCYAN, BDRCYAN,
    ),
]
RW, RH, RGX, RGY = 3.78, 1.88, 0.52, 3.78
RGAP = 0.22
for i, (title, steps, col, fill, bdr) in enumerate(retrievers):
    rx = RGX + i * (RW + RGAP)
    flow_box(sl3, rx, RGY, RW, RH, title, steps,
             title_color=col, fill=fill, border=bdr)

# ── ISD AGENT LOOP BAND ──────────────────────────────────
rect(sl3, 0.42, 5.88, 12.50, 1.38, fill=RGBColor(0x06,0x0E,0x06), border=BDRGREEN)
txt(sl3, "ISD AGENT LOOP  ·  ONE JOB PER (DOCUMENT × COLUMN) PAIR",
    0.60, 5.93, 9, 0.24, size=7.5, bold=True, color=BDRGREEN)

agent_steps = [
    ("Decompose",    "Breaks prompt into\nsub-questions +\ntarget sections",    VIOLET, DKVIOLET, BDRVIOLET),
    ("Gather",       "Calls retriever per\nsub-question,\npools evidence",       BRBLUE, DKBLUE,   BDRBLUE),
    ("Draft",        "LLM writes answer\nfrom evidence only,\ncites every claim",SOFTGRN,DKSFGRN,  BDRSFGRN),
    ("Verify",       "Re-retrieves claims,\nchecks supported /\ncontradicted",   ORANGE, DKORANGE,  BDRORANGE),
    ("Cell Result",  "Answer + citations\nconfidence +\nfull trace",             CYAN,   DKCYAN,    BDRCYAN),
]
AW, AH, AGX, AGY = 2.14, 1.02, 0.52, 6.14
AGAP = 0.24
for i, (title, sub, col, fill, bdr) in enumerate(agent_steps):
    ax = AGX + i * (AW + AGAP)
    rect(sl3, ax, AGY, AW, AH, fill=fill, border=bdr)
    txt(sl3, title, ax+0.13, AGY+0.09, AW-0.24, 0.26, size=9, bold=True, color=col)
    txt(sl3, sub,   ax+0.13, AGY+0.38, AW-0.24, 0.58, size=8, color=LGREY)
    if i < 4:
        arrow_right(sl3, ax + AW, AGY + AH/2,
                    ax + AW + AGAP, col=BDRGREEN, lw=Pt(1.0))

# FastAPI + React at far right
rect(sl3, 11.22, 6.14, 1.60, 1.02, fill=DKBLUE, border=BDRBLUE)
txt(sl3, "FastAPI + SSE", 11.35, 6.22, 1.35, 0.26, size=8, bold=True, color=BRBLUE)
txt(sl3, "Streams cell state\nto React UI",
    11.35, 6.52, 1.35, 0.50, size=7.5, color=GREY)
arrow_right(sl3,
    AGX + 5*(AW+AGAP) - AGAP, AGY + AH/2,
    11.22, col=BDRGREEN, lw=Pt(1.0))

# ── save ─────────────────────────────────────────────────
out = "/Users/anup.roy/Downloads/Hebbia POC/docs/hebbia-poc-slides.pptx"
prs.save(out)
print(f"Saved: {out}")
